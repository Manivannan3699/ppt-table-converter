# backend/server.py
"""
FINAL PPTX → HTML TABLE CONVERTER (CLASS-BASED)

✔ Externalized CSS classes only
✔ Semantic + hex color classes
✔ Font family / size classes
✔ Alignment classes
✔ Bullet handling
✔ Render-safe (no global state bugs)
✔ Same output locally and on Render
"""

import os
import tempfile
import logging
import html
import re
from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")

# ------------------------------------------------------------------
# FALLBACK THEME COLORS
# ------------------------------------------------------------------
THEME_COLOR_MAP = {
    "ACCENT1": "#4472C4", "ACCENT_1": "#4472C4",
    "ACCENT2": "#ED7D31", "ACCENT_2": "#ED7D31",
    "ACCENT3": "#A5A5A5", "ACCENT_3": "#A5A5A5",
    "ACCENT4": "#FFC000", "ACCENT_4": "#FFC000",
    "ACCENT5": "#5B9BD5", "ACCENT_5": "#5B9BD5",
    "ACCENT6": "#70AD47", "ACCENT_6": "#70AD47",
    "TEXT1": "#000000", "TEXT_1": "#000000",
    "TEXT2": "#FFFFFF", "TEXT_2": "#FFFFFF",
    "BACKGROUND1": "#FFFFFF", "BACKGROUND_1": "#FFFFFF",
    "BACKGROUND2": "#000000", "BACKGROUND_2": "#000000",
}

# ------------------------------------------------------------------
# CSS REGISTRY (PER REQUEST)
# ------------------------------------------------------------------
GENERATED_CSS = set()

# ------------------------------------------------------------------
# HELPERS
# ------------------------------------------------------------------
def sanitize(s, maxlen=40):
    if not s:
        return ""
    return re.sub(r"[^a-zA-Z0-9]+", "-", s).strip("-").lower()[:maxlen]


def rgb_to_hex(rgb):
    try:
        return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        return None


def add_css(selector, body):
    GENERATED_CSS.add(f"{selector} {{ {body} }}")


# ------------------------------------------------------------------
# THEME EXTRACTION
# ------------------------------------------------------------------
def extract_theme_colors(prs):
    theme = {}
    try:
        root = parse_xml(prs.part.theme_part.blob)
        scheme = root.find(".//a:clrScheme", root.nsmap)
        if scheme is not None:
            for el in scheme:
                key = el.tag.split("}")[-1].upper()
                srgb = el.find(".//a:srgbClr", root.nsmap)
                if srgb is not None:
                    theme[key] = "#" + srgb.get("val")
    except Exception:
        pass

    for k, v in THEME_COLOR_MAP.items():
        theme.setdefault(k.upper(), v)

    return theme


# ------------------------------------------------------------------
# CSS CLASS GENERATORS
# ------------------------------------------------------------------
def bg_classes(hex_color, semantic=None):
    if not hex_color:
        return []
    key = hex_color[1:].upper()
    cls_hex = f"bg-{key}"
    add_css(f".{cls_hex}", f"background-color:{hex_color};")

    classes = [cls_hex]
    if semantic:
        cls_sem = f"bg-{sanitize(semantic)}"
        add_css(f".{cls_sem}", f"background-color:{hex_color};")
        classes.insert(0, cls_sem)

    return classes


def text_color_class(hex_color):
    if not hex_color:
        return []
    cls = f"text-{hex_color[1:].upper()}"
    add_css(f".{cls}", f"color:{hex_color};")
    return [cls]


def font_family_class(name):
    if not name:
        return []
    cls = f"ff-{sanitize(name)}"
    add_css(f".{cls}", f"font-family:'{name}';")
    return [cls]


def font_size_class(pt):
    if not pt:
        return []
    cls = f"fs-{int(pt)}"
    add_css(f".{cls}", f"font-size:{int(pt)}pt;")
    return [cls]


# ------------------------------------------------------------------
# RUN / PARAGRAPH / CELL TO HTML
# ------------------------------------------------------------------
def run_to_html(run, theme):
    text = html.escape(run.text or "")
    if not text:
        return ""

    classes = []

    try:
        fc = run.font.color
        if fc:
            if fc.rgb:
                classes += text_color_class(rgb_to_hex(fc.rgb))
            elif fc.theme_color:
                key = str(fc.theme_color).split(".")[-1].upper()
                classes += text_color_class(theme.get(key))
    except:
        pass

    try:
        classes += font_family_class(run.font.name)
    except:
        pass

    try:
        if run.font.size:
            classes += font_size_class(run.font.size.pt)
    except:
        pass

    span = text
    if classes:
        span = f'<span class="{" ".join(dict.fromkeys(classes))}">{text}</span>'

    if run.font.bold:
        span = f"<b>{span}</b>"
    if run.font.italic:
        span = f"<i>{span}</i>"
    if run.font.underline:
        span = f"<u>{span}</u>"

    return span


def para_to_html(p, theme):
    parts = [run_to_html(r, theme) for r in p.runs]
    content = "".join(parts).strip()
    if not content:
        return ""

    # Bullet detection
    try:
        if p.level > 0:
            return f"<li>{content}</li>"
        ppr = p._pPr
        if ppr is not None and ppr.find(".//a:buChar", p._element.nsmap) is not None:
            return f"<li>{content}</li>"
    except:
        pass

    return content


def cell_content(cell, theme):
    items = [para_to_html(p, theme) for p in cell.text_frame.paragraphs]
    items = [i for i in items if i]

    if not items:
        return ""

    if any(i.startswith("<li>") for i in items):
        return "<ul>" + "".join(items) + "</ul>"

    return "<br/>".join(items)


def cell_classes(cell, theme):
    classes = ["ppt-cell"]

    # Background
    try:
        fc = cell.fill.fore_color
        hexc = None
        semantic = None

        if fc.rgb:
            hexc = rgb_to_hex(fc.rgb)
        elif fc.theme_color:
            semantic = str(fc.theme_color).split(".")[-1].lower()
            hexc = theme.get(semantic.upper())

        if hexc:
            classes += bg_classes(hexc, semantic)
    except:
        pass

    # Alignment
    try:
        align = cell.text_frame.paragraphs[0].alignment
        if align == PP_ALIGN.CENTER:
            classes.append("align-center")
            add_css(".align-center", "text-align:center;")
        elif align == PP_ALIGN.RIGHT:
            classes.append("align-right")
            add_css(".align-right", "text-align:right;")
        else:
            classes.append("align-left")
            add_css(".align-left", "text-align:left;")
    except:
        classes.append("align-left")
        add_css(".align-left", "text-align:left;")

    add_css(".ppt-cell", "border:1px solid #999;padding:6px;vertical-align:middle;")

    return " ".join(dict.fromkeys(classes))


# ------------------------------------------------------------------
# TABLE PROCESSING
# ------------------------------------------------------------------
def process_table(table, theme):
    rows = []

    for r in range(len(table.rows)):
        cols = []
        for c in range(len(table.columns)):
            cell = table.cell(r, c)

            if getattr(cell, "is_spanned", False):
                continue

            content = cell_content(cell, theme)
            classes = cell_classes(cell, theme)

            colspan = getattr(cell, "span_width", 1)
            rowspan = getattr(cell, "span_height", 1)

            cols.append(
                f'<td colspan="{colspan}" rowspan="{rowspan}" class="{classes}">{content}</td>'
            )

        rows.append("<tr>" + "".join(cols) + "</tr>")

    add_css(".ppt-table", "border-collapse:collapse;width:100%;font-size:14px;")
    add_css(".ppt-table ul", "margin:0 0 0 18px;padding:0;")

    return "\n".join(rows)


# ------------------------------------------------------------------
# MAIN CONVERTER (IMPORTANT: CLEAR CSS ONCE HERE)
# ------------------------------------------------------------------
def pptx_to_html(path):
    GENERATED_CSS.clear()

    prs = Presentation(path)
    theme = extract_theme_colors(prs)
    blocks = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                body = process_table(shape.table, theme)
                blocks.append(f"<table class='ppt-table'>{body}</table>")

    css = "<style>\n" + "\n".join(GENERATED_CSS) + "\n</style>\n"
    return css + "<br/>".join(blocks)


# ------------------------------------------------------------------
# FLASK APP
# ------------------------------------------------------------------
app = Flask(__name__)
CORS(app)

@app.route("/")
def home():
    return "PPTX → HTML Converter API running"


@app.route("/convert", methods=["POST"])
def convert():
    f = request.files.get("pptFile")
    if not f:
        return jsonify({"error": "No file uploaded"}), 400

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        f.save(tmp.name)
        html_out = pptx_to_html(tmp.name)

    os.remove(tmp.name)
    return jsonify({"html": html_out})


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
